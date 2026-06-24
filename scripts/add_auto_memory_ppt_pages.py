from copy import deepcopy
from datetime import datetime
from pathlib import Path
from shutil import copy2

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt


PPT_PATH = Path(r"C:\Users\dd\Documents\PowerShell\Claude_Code_0604.pptx")


def rgb(hexstr: str) -> RGBColor:
    return RGBColor(int(hexstr[0:2], 16), int(hexstr[2:4], 16), int(hexstr[4:6], 16))


def clone_slide(prs: Presentation, source_slide):
    blank = prs.slides.add_slide(source_slide.slide_layout)
    for shape in list(blank.shapes):
        blank.shapes._spTree.remove(shape.element)
    for shape in source_slide.shapes:
        blank.shapes._spTree.insert_element_before(deepcopy(shape.element), "p:extLst")
    return blank


def set_text(shape, text: str, size=None, bold=None, color=None):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    if size:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if color:
        run.font.color.rgb = color


def set_multiline(shape, title: str, body: str, title_size=13, body_size=9):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(title_size)
    r.font.bold = True
    r.font.color.rgb = rgb("111111")
    if body:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = body
        r2.font.name = "Microsoft YaHei"
        r2.font.size = Pt(body_size)
        r2.font.color.rgb = rgb("555555")


pages = [
    {
        "title": "Auto Memory：文件化的长期记忆",
        "subtitle": "不是数据库或向量库，而是项目级 Markdown 目录 + MEMORY.md 索引。",
        "left_title": "存储结构",
        "steps": [
            ("项目级目录", "~/.claude/projects/<repo>/memory/；同一 Git 仓库的 worktree 共享。"),
            ("MEMORY.md 入口", "每次会话加载；只放链接索引，不直接写长内容。"),
            ("独立记忆文件", "每条记忆一个 Markdown 文件，带 name / description / type frontmatter。"),
            ("路径安全", "autoMemoryDirectory 只信任 policy/local/user；故意排除 projectSettings。"),
        ],
        "callout": "getAutoMemPath() → getAutoMemEntrypoint() → MEMORY.md",
        "caption": "入口索引短而稳，细节散落到主题文件",
        "right_title": "目录与文件",
        "headers": ("文件", "职责"),
        "rows": [
            ("MEMORY.md", "入口索引；200 行 / 25KB 上限"),
            ("user_*.md", "用户角色、偏好、背景"),
            ("feedback_*.md", "用户纠正或确认过的行为"),
            ("project_*.md", "非代码可推导的项目上下文"),
        ],
        "bottom": "关键点：记忆是可读写文件系统，不是隐藏黑盒。",
    },
    {
        "title": "Auto Memory：什么值得被保存",
        "subtitle": "只保存未来会话有用、且无法从当前仓库状态直接推导的信息。",
        "left_title": "四类型分类法",
        "steps": [
            ("user", "用户是谁、技术背景、个人协作偏好。"),
            ("feedback", "对 AI 行为的纠正与正向确认；成功和失败都要记录。"),
            ("project", "人、时间、动机、约束；变化快，所以必须带 Why / How to apply。"),
            ("reference", "外部系统指针：Linear 项目、仪表盘、常用服务入口。"),
        ],
        "callout": "memoryTypes.ts：<when_to_save> / <how_to_use> / <body_structure>",
        "caption": "代码能 grep 到的事实，不应该写成记忆",
        "right_title": "不要保存什么",
        "headers": ("不要记", "为什么"),
        "rows": [
            ("文件路径清单", "仓库可实时搜索"),
            ("代码架构事实", "会过期，可从源码推导"),
            ("一次性任务进度", "应该用 Todo / Plan"),
            ("敏感信息", "尤其不能写入团队记忆"),
        ],
        "bottom": "记忆不是摘要；它是跨会话的行为和上下文约束。",
    },
    {
        "title": "Auto Memory：召回与注入链路",
        "subtitle": "每次对话不盲目塞全部记忆，而是索引常驻 + 相关记忆按需召回。",
        "left_title": "运行时链路",
        "steps": [
            ("系统提示注入", "loadMemoryPrompt() 通过 systemPromptSection('memory') 装配行为规则。"),
            ("索引进入上下文", "MEMORY.md 作为入口上下文加载，让模型知道可用记忆有哪些。"),
            ("Sonnet 侧查询", "findRelevantMemories() 扫描 frontmatter，最多挑选少量相关文件。"),
            ("读取前防漂移", "Before recommending from memory：路径先检查，函数/flag 先 grep。"),
        ],
        "callout": "用户消息 → manifest → selectRelevantMemories() → Read 相关文件",
        "caption": "召回是相关性选择，不是全量灌上下文",
        "right_title": "缓存友好设计",
        "headers": ("层", "作用"),
        "rows": [
            ("memory section", "稳定行为规则，可缓存"),
            ("MEMORY.md", "短索引，控制上下文成本"),
            ("relevant files", "按需读取，避免噪声"),
            ("alreadySurfaced", "避免重复召回同一文件"),
        ],
        "bottom": "目标：少量、高信号、可验证地进入当前会话。",
    },
    {
        "title": "Auto Memory：后台提炼与上下文压缩",
        "subtitle": "记忆系统不只在主会话里工作，还会在后台整理、梦境合并、压缩续接。",
        "left_title": "后台机制",
        "steps": [
            ("Extract Memories", "回合结束后 fork 子代理，从最近消息里抽取可保存信息。"),
            ("hasMemoryWritesSince", "主代理已经写过记忆时，后台抽取跳过对应范围，避免重复。"),
            ("Dream / AutoDream", "对 memory files 和 transcript 做反思整理，合并为更耐用主题。"),
            ("Session Memory Compact", "压缩时可优先用 session memory 续接，而不是再调用摘要模型。"),
        ],
        "callout": "extractMemories → memory files → /dream consolidation → compact resume",
        "caption": "记忆在会话外继续生长，但仍受 prompt 规则约束",
        "right_title": "主会话之外",
        "headers": ("模块", "职责"),
        "rows": [
            ("extractMemories", "后台抽取新偏好"),
            ("consolidationPrompt", "梦境式整理长期记忆"),
            ("Team Memory", "私有 / 团队目录分流"),
            ("SM Compact", "用会话记忆降低压缩成本"),
        ],
        "bottom": "AutoMemory = 写入规则 + 召回策略 + 后台整理。",
    },
]


backup_path = PPT_PATH.with_name(
    PPT_PATH.stem
    + "_backup_before_auto_memory_pages_"
    + datetime.now().strftime("%Y%m%d_%H%M%S")
    + PPT_PATH.suffix
)
copy2(PPT_PATH, backup_path)

prs = Presentation(str(PPT_PATH))
source = prs.slides[23]  # current Auto Memory overview slide
new_slides = [clone_slide(prs, source) for _ in pages]

sld_id_lst = prs.slides._sldIdLst
ids = list(sld_id_lst)
appended = ids[-len(pages) :]
for el in appended:
    sld_id_lst.remove(el)
insert_at = 24  # after current slide 24
for offset, el in enumerate(appended):
    sld_id_lst.insert(insert_at + offset, el)

targets = [prs.slides[i] for i in range(24, 24 + len(pages))]

for slide, page in zip(targets, pages):
    shapes = list(slide.shapes)
    set_text(shapes[0], page["title"], size=22, bold=True, color=rgb("3A1D07"))
    set_text(shapes[1], page["subtitle"], size=12, bold=False, color=rgb("555555"))
    set_text(shapes[3], page["left_title"], size=17, bold=True, color=rgb("3A1D07"))

    number_shapes = [shapes[i] for i in [4, 6, 8, 10]]
    step_text_shapes = [shapes[i] for i in [5, 7, 9, 11]]
    for idx, (num_shape, text_shape, (step_title, step_body)) in enumerate(
        zip(number_shapes, step_text_shapes, page["steps"]), 1
    ):
        set_text(num_shape, str(idx), size=8, bold=True, color=rgb("FFFFFF"))
        set_multiline(text_shape, step_title, step_body, title_size=12.5, body_size=8.5)

    set_text(shapes[12], page["callout"], size=10.5, bold=True, color=rgb("111111"))
    set_text(shapes[13], page["caption"], size=8.5, bold=False, color=rgb("555555"))
    set_text(shapes[14], page["right_title"], size=15.5, bold=True, color=rgb("3A1D07"))

    set_text(shapes[16], page["headers"][0], size=10.5, bold=True, color=rgb("FFFFFF"))
    set_text(shapes[17], page["headers"][1], size=10.5, bold=True, color=rgb("FFFFFF"))

    row_shapes = [(18, 19), (20, 21), (22, 23), (24, 25)]
    for (left_idx, right_idx), (left_text, right_text) in zip(row_shapes, page["rows"]):
        set_text(shapes[left_idx], left_text, size=9.2, bold=False, color=rgb("111111"))
        set_text(shapes[right_idx], right_text, size=9.2, bold=False, color=rgb("111111"))

    set_text(shapes[26], page["bottom"], size=10.2, bold=True, color=rgb("FFFFFF"))

prs.save(str(PPT_PATH))
print(f"backup {backup_path}")
print(f"saved {PPT_PATH}")
print(f"slides {len(prs.slides)}")
